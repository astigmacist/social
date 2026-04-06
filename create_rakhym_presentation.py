from __future__ import annotations

import os
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


BASE_DIR = Path("/Users/erbolsadibekov/Desktop/social")
OUTPUT = Path("/Users/erbolsadibekov/Desktop/Rakhym_presentation.pptx")
LOGO = BASE_DIR / "logo.png"


TITLE_COLOR = RGBColor(26, 43, 76)
TEXT_COLOR = RGBColor(62, 77, 102)
MUTED_COLOR = RGBColor(111, 126, 150)
PANEL_COLOR = RGBColor(255, 255, 255)
PANEL_DARK = RGBColor(245, 249, 253)
ACCENT = RGBColor(21, 94, 117)
ACCENT_2 = RGBColor(180, 143, 74)
ACCENT_3 = RGBColor(59, 130, 246)
BG = RGBColor(248, 250, 252)
BG_DARK = RGBColor(241, 245, 249)
BORDER = RGBColor(197, 210, 226)


def set_background(slide, variant: int = 1) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG if variant == 1 else BG_DARK

    top = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.45)
    )
    top.fill.solid()
    top.fill.fore_color.rgb = RGBColor(226, 236, 248)
    top.line.fill.background()

    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, Inches(0.45), prs.slide_width, Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    for left, top_pos, size, color in [
        (Inches(11.8), Inches(0.15), Inches(0.16), RGBColor(219, 234, 254)),
        (Inches(12.05), Inches(0.42), Inches(0.24), RGBColor(224, 231, 255)),
        (Inches(12.3), Inches(0.14), Inches(0.16), RGBColor(209, 250, 229)),
        (Inches(0.15), Inches(6.65), Inches(0.18), RGBColor(226, 232, 240)),
        (Inches(0.4), Inches(6.4), Inches(0.26), RGBColor(219, 234, 254)),
    ]:
        deco = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top_pos, size, size
        )
        deco.fill.solid()
        deco.fill.fore_color.rgb = color
        deco.fill.transparency = 0.08
        deco.line.fill.background()

    glow = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL, Inches(11.1), Inches(0.5), Inches(1.8), Inches(1.8)
    )
    glow.fill.solid()
    glow.fill.fore_color.rgb = RGBColor(191, 219, 254)
    glow.fill.transparency = 0.84
    glow.line.fill.background()


def add_logo(slide, left=Inches(0.22), top=Inches(0.08), width=Inches(0.55)) -> None:
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), left, top, width=width)


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.65), Inches(0.6), Inches(8.8), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Arial"
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = TITLE_COLOR

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.68), Inches(1.05), Inches(10.0), Inches(0.4))
        tf2 = sub_box.text_frame
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.name = "Arial"
        run2.font.size = Pt(10.5)
        run2.font.color.rgb = MUTED_COLOR


def add_slide_number(slide, num: int) -> None:
    box = slide.shapes.add_textbox(Inches(12.45), Inches(6.95), Inches(0.5), Inches(0.2))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = str(num)
    r.font.name = "Arial"
    r.font.size = Pt(9)
    r.font.color.rgb = MUTED_COLOR


def add_panel(slide, left, top, width, height, title=None, body=None, fill_color=PANEL_COLOR):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.fill.transparency = 0.0
    shape.line.color.rgb = BORDER
    shape.line.transparency = 0.0

    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.18)
    tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.10)
    tf.margin_bottom = Inches(0.08)

    if title:
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = title
        r.font.name = "Arial"
        r.font.size = Pt(14)
        r.font.bold = True
        r.font.color.rgb = TITLE_COLOR

    if body:
        if isinstance(body, str):
            body = [body]
        for idx, line in enumerate(body):
            p = tf.add_paragraph() if (title or idx > 0) else tf.paragraphs[0]
            if not title and idx == 0:
                p = tf.paragraphs[0]
            p.text = line
            p.level = 0
            p.font.name = "Arial"
            p.font.size = Pt(10.5)
            p.font.color.rgb = TEXT_COLOR
            p.space_after = Pt(2)
    return shape


def add_bullet_box(slide, left, top, width, height, title, bullets, fill_color=PANEL_DARK):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.fill.transparency = 0.0
    shape.line.color.rgb = BORDER

    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.18)
    tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.1)

    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Arial"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT

    for bullet in bullets:
        bp = tf.add_paragraph()
        bp.text = bullet
        bp.level = 0
        bp.bullet = True
        bp.font.name = "Arial"
        bp.font.size = Pt(10.5)
        bp.font.color.rgb = TEXT_COLOR
        bp.space_after = Pt(1)
    return shape


def add_stat_card(slide, left, top, width, height, number, label, color=ACCENT):
    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    card.fill.solid()
    card.fill.fore_color.rgb = PANEL_COLOR
    card.fill.transparency = 0.0
    card.line.color.rgb = BORDER

    tf = card.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = str(number)
    r.font.name = "Arial"
    r.font.size = Pt(26)
    r.font.bold = True
    r.font.color.rgb = color

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = label
    r2.font.name = "Arial"
    r2.font.size = Pt(10)
    r2.font.color.rgb = TEXT_COLOR


def add_cover_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, variant=1)
    add_logo(slide)

    badge = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.72), Inches(0.92), Inches(3.2), Inches(0.42)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = RGBColor(230, 240, 250)
    badge.line.color.rgb = BORDER
    badge.text_frame.text = "Әлеуметтік жұмыс олимпиадасына арналған жоба"
    badge.text_frame.paragraphs[0].font.name = "Arial"
    badge.text_frame.paragraphs[0].font.size = Pt(11)
    badge.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR

    title_box = slide.shapes.add_textbox(Inches(0.72), Inches(1.55), Inches(5.3), Inches(2.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "RAKHYM"
    r.font.name = "Arial"
    r.font.size = Pt(29)
    r.font.bold = True
    r.font.color.rgb = ACCENT

    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = "Әлеуметтік қолдауға арналған\nцифрлық веб-платформа"
    r2.font.name = "Arial"
    r2.font.size = Pt(22)
    r2.font.bold = True
    r2.font.color.rgb = TITLE_COLOR

    p3 = tf.add_paragraph()
    r3 = p3.add_run()
    r3.text = (
        "Жоба азаматтарға жеңілдіктерді табуға, онлайн өтініш беруге,\n"
        "мәртебені бақылауға және маман/ерікті сүйемелдеуін алуға көмектеседі."
    )
    r3.font.name = "Arial"
    r3.font.size = Pt(12)
    r3.font.color.rgb = TEXT_COLOR

    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.15), Inches(1.25), Inches(4.15), Inches(3.95)
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = PANEL_COLOR
    panel.fill.transparency = 0.0
    panel.line.color.rgb = BORDER

    ptf = panel.text_frame
    ptf.clear()
    items = [
        ("Жоба түрі", "Social Tech / цифрлық әлеуметтік сервис"),
        ("Форматы", "Веб-платформа + әкімшілік панель + жеке кабинет"),
        ("Тілдер", "Қазақша және орысша"),
        ("Күйі", "Негізгі модульдері бар жұмыс істейтін прототип"),
    ]
    for idx, (label, value) in enumerate(items):
        p = ptf.paragraphs[0] if idx == 0 else ptf.add_paragraph()
        p.text = f"{label}: {value}"
        p.font.name = "Arial"
        p.font.size = Pt(12 if idx == 0 else 11)
        p.font.color.rgb = TITLE_COLOR if idx == 0 else TEXT_COLOR
        if idx == 0:
            p.font.bold = True

    bubble = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL, Inches(8.9), Inches(5.5), Inches(2.45), Inches(0.8)
    )
    bubble.fill.solid()
    bubble.fill.fore_color.rgb = RGBColor(254, 249, 195)
    bubble.fill.transparency = 0.0
    bubble.line.color.rgb = RGBColor(234, 179, 8)
    bubble.text_frame.text = "Түсінікті. Қолжетімді. Адамға жақын."
    bubble.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    bubble.text_frame.paragraphs[0].font.name = "Arial"
    bubble.text_frame.paragraphs[0].font.size = Pt(11)
    bubble.text_frame.paragraphs[0].font.bold = True
    bubble.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR

    palette_title = slide.shapes.add_textbox(Inches(0.82), Inches(5.55), Inches(1.7), Inches(0.25))
    pp = palette_title.text_frame.paragraphs[0]
    rr = pp.add_run()
    rr.text = "Цветограмма"
    rr.font.name = "Arial"
    rr.font.size = Pt(10)
    rr.font.bold = True
    rr.font.color.rgb = MUTED_COLOR

    for idx, color in enumerate([ACCENT, ACCENT_3, ACCENT_2, RGBColor(30, 41, 59)]):
        chip = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            Inches(0.85 + idx * 0.35),
            Inches(5.88),
            Inches(0.18),
            Inches(0.18),
        )
        chip.fill.solid()
        chip.fill.fore_color.rgb = color
        chip.line.fill.background()

    add_slide_number(slide, 1)


def add_relevance_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, variant=1)
    add_logo(slide)
    add_title(slide, "Тақырыптың өзектілігі", "Неліктен Rakhym сияқты платформа қажет?")

    add_bullet_box(
        slide,
        Inches(0.7),
        Inches(1.55),
        Inches(5.6),
        Inches(4.6),
        "Қазіргі әлеуметтік ортадағы негізгі қиындықтар",
        [
            "Көп адам қандай жеңілдіктер мен жәрдемақылар бар екенін білмейді.",
            "Қажетті құжаттар мен келесі қадамдарды түсіну қиын.",
            "Ақпарат, өтініш беру және сүйемелдеу әртүрлі жерде орналасқан.",
            "Осал топтарға түсінікті әрі адамға жақын цифрлық сервис жетіспейді.",
        ],
    )

    add_panel(
        slide,
        Inches(6.6),
        Inches(1.7),
        Inches(5.7),
        Inches(1.2),
        "Rakhym идеясы",
        "Әлеуметтік көмекті бір жерде шоғырландырып, азаматты ақпараттан өтінішке дейін толық цифрлық маршрутпен сүйемелдеу.",
    )
    add_panel(
        slide,
        Inches(6.6),
        Inches(3.05),
        Inches(2.7),
        Inches(1.35),
        "Кімдерге маңызды?",
        ["Аз қамтылған отбасылар", "Қарттар", "Жұмыссыздар", "Мүгедектігі бар адамдар"],
        fill_color=PANEL_DARK,
    )
    add_panel(
        slide,
        Inches(9.55),
        Inches(3.05),
        Inches(2.75),
        Inches(1.35),
        "Практикалық мәселе",
        ["Көмекке құқығы бар адамдар оны толық пайдалана алмайды."],
        fill_color=PANEL_DARK,
    )
    add_panel(
        slide,
        Inches(6.6),
        Inches(4.6),
        Inches(5.7),
        Inches(1.55),
        "Нәтижесінде",
        [
            "Көмекке жүгіну кедергісі жоғарылайды, ал азамат үшін процесс түсініксіз болып қалады.",
            "Rakhym осы кедергіні азайтуға бағытталған."
        ],
        fill_color=PANEL_DARK,
    )
    add_slide_number(slide, 2)


def add_problem_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, variant=2)
    add_logo(slide)
    add_title(slide, "Rakhym шешетін негізгі мәселе", "Платформа бірнеше байланысты проблеманы бір уақытта жабады")

    items = [
        ("01", "Ақпаратқа қол жеткізу қиындығы", "Азаматтар қандай қолдау барын және оған сай келетін-келмейтінін түсінбейді."),
        ("02", "Өтініш беру қиындығы", "Форманы толтыру, құжат жинау, келесі қадамды түсіну күрделі."),
        ("03", "Бірыңғай кіру нүктесінің болмауы", "Ақпарат, көмек, еріктілер және сүйемелдеу әртүрлі арналарға бөлінген."),
        ("04", "Адамға жақын сүйемелдеудің жетіспеуі", "Адамға тек мәлімет емес, түсіндіру мен нақты бағыт-бағдар қажет."),
    ]
    positions = [
        (Inches(0.75), Inches(1.6)),
        (Inches(6.65), Inches(1.6)),
        (Inches(0.75), Inches(4.0)),
        (Inches(6.65), Inches(4.0)),
    ]
    for (num, title, body), (left, top) in zip(items, positions):
        box = add_panel(slide, left, top, Inches(5.25), Inches(1.8), fill_color=PANEL_DARK)
        tf = box.text_frame
        tf.clear()
        p0 = tf.paragraphs[0]
        r0 = p0.add_run()
        r0.text = num
        r0.font.name = "Arial"
        r0.font.size = Pt(18)
        r0.font.bold = True
        r0.font.color.rgb = ACCENT

        p1 = tf.add_paragraph()
        p1.text = title
        p1.font.name = "Arial"
        p1.font.size = Pt(13)
        p1.font.bold = True
        p1.font.color.rgb = TITLE_COLOR

        p2 = tf.add_paragraph()
        p2.text = body
        p2.font.name = "Arial"
        p2.font.size = Pt(10.3)
        p2.font.color.rgb = TEXT_COLOR
    add_slide_number(slide, 3)


def add_goal_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, variant=1)
    add_logo(slide)
    add_title(slide, "Жобаның мақсаты мен міндеттері")

    add_panel(
        slide,
        Inches(0.75),
        Inches(1.55),
        Inches(11.6),
        Inches(1.0),
        "Басты мақсат",
        "Азаматтарға әлеуметтік қолдау шараларын тез табуға, өтініш беруге және оны түсінікті форматта сүйемелдеуге көмектесетін ыңғайлы цифрлық платформа құру.",
        fill_color=PANEL_DARK,
    )

    add_bullet_box(
        slide,
        Inches(0.75),
        Inches(2.85),
        Inches(5.7),
        Inches(3.15),
        "Негізгі міндеттер",
        [
            "Жеңілдіктер мен жәрдемақылар туралы ақпаратты түсінікті ету",
            "Өтініш беру процесін жеңілдету",
            "Пайдаланушы үшін бірыңғай цифрлық маршрут құру",
            "Азаматтар, еріктілер және мамандар байланысын ұйымдастыру",
        ],
    )
    add_bullet_box(
        slide,
        Inches(6.65),
        Inches(2.85),
        Inches(5.7),
        Inches(3.15),
        "Қосымша міндеттер",
        [
            "Көмекке жүгіну кедергісін азайту",
            "Әлеуметтік қызметтердің цифрлық қолжетімділігін арттыру",
            "Әкімшілік басқаруды жеңілдету",
            "AI-құралдарын енгізуге негіз дайындау",
        ],
        fill_color=PANEL_DARK,
    )
    add_slide_number(slide, 4)


def add_audience_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, variant=2)
    add_logo(slide)
    add_title(slide, "Мақсатты аудитория", "Rakhym бірнеше пайдаланушы тобымен жұмыс істейді")

    add_bullet_box(
        slide,
        Inches(0.75),
        Inches(1.6),
        Inches(5.75),
        Inches(4.6),
        "Негізгі аудитория",
        [
            "Әлеуметтік көмекке мұқтаж азаматтар",
            "Балалы отбасылар",
            "Аз қамтылған отбасылар",
            "Қарт адамдар",
            "Мүгедектігі бар адамдар",
            "Жұмыссыздар",
        ],
    )
    add_bullet_box(
        slide,
        Inches(6.7),
        Inches(1.6),
        Inches(5.55),
        Inches(4.6),
        "Қосымша аудитория",
        [
            "Еріктілер",
            "Әлеуметтік қызметкерлер",
            "Өтініштерді сүйемелдеу мамандары",
            "Әкімшілер",
            "ҮЕҰ және серіктес ұйымдар",
        ],
        fill_color=PANEL_DARK,
    )
    add_slide_number(slide, 5)


def add_solution_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, variant=1)
    add_logo(slide)
    add_title(slide, "Rakhym платформасы не ұсынады?", "Бір шешім ішінде бірнеше толыққанды модуль біріктірілген")

    cards = [
        ("Жеңілдіктер каталогы", "Санаттар, карточкалар, ресми дереккөздер"),
        ("Онлайн өтініш", "Форма, өтініш нөмірі, мәртебе"),
        ("Жеке кабинет", "Профиль, өтініштер, баптаулар"),
        ("Еріктілер модулі", "Тіркеу, бағыт таңдау, мәртебе"),
        ("Жаңалықтар блогы", "Өзекті хабарландырулар мен өзгерістер"),
        ("Цифрлық көмекші", "Типтік сұрақтарға жылдам жауап"),
    ]
    xs = [Inches(0.8), Inches(4.35), Inches(7.9)]
    ys = [Inches(1.8), Inches(4.1)]
    idx = 0
    for y in ys:
        for x in xs:
            title, body = cards[idx]
            add_panel(slide, x, y, Inches(3.15), Inches(1.8), title, body)
            idx += 1
    add_slide_number(slide, 6)


def add_journey_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, variant=2)
    add_logo(slide)
    add_title(slide, "Пайдаланушы жолы", "Rakhym логикасы: сұрақтан бастап сүйемелдеуге дейін")

    steps = [
        "1. Сайтқа кіреді",
        "2. Жеңілдікті немесе көмек түрін іздейді",
        "3. Онлайн өтініш береді",
        "4. Өтініш нөмірін алады",
        "5. Жеке кабинеттен мәртебені бақылайды",
        "6. Қажет болса маман/ерікті қосылады",
    ]
    left = Inches(0.9)
    top = Inches(2.2)
    width = Inches(1.8)
    gap = Inches(0.22)
    shapes = []
    for i, step in enumerate(steps):
        shp = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            left + i * (width + gap),
            top,
            width,
            Inches(1.25),
        )
        shp.fill.solid()
        shp.fill.fore_color.rgb = PANEL_COLOR
        shp.line.color.rgb = BORDER
        shp.text_frame.text = step
        shp.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        shp.text_frame.paragraphs[0].font.name = "Arial"
        shp.text_frame.paragraphs[0].font.size = Pt(10.2)
        shp.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
        shp.text_frame.paragraphs[0].font.bold = True
        shapes.append(shp)
    for i in range(len(shapes) - 1):
        x1 = shapes[i].left + shapes[i].width
        y1 = shapes[i].top + shapes[i].height // 2
        x2 = shapes[i + 1].left
        y2 = shapes[i + 1].top + shapes[i + 1].height // 2
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
        conn.line.color.rgb = ACCENT_3
        conn.line.width = Pt(2)

    add_panel(
        slide,
        Inches(1.2),
        Inches(4.4),
        Inches(10.8),
        Inches(1.35),
        "Негізгі идея",
        "Rakhym жай анықтамалық емес: ол ақпаратты, өтінішті, бақылауды және адами сүйемелдеуді бір цифрлық өнімге біріктіреді.",
        fill_color=PANEL_DARK,
    )
    add_slide_number(slide, 7)


def add_features_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, variant=1)
    add_logo(slide)
    add_title(slide, "Қазірдің өзінде іске асқан функциялар", "Жоба идея деңгейінде емес, жұмыс істейтін прототип деңгейінде")

    features = [
        "Басты бет және жылдам навигация",
        "Санаттарға бөлінген жеңілдіктер каталогы",
        "Онлайн өтініш беру формасы",
        "Пайдаланушының жеке кабинеті",
        "Тіркеу және авторизация",
        "Еріктілерді тіркеу модулі",
        "Жаңалықтар және хабарландырулар",
        "Әкімшілік панель және өтініштерді басқару",
    ]
    for i, feature in enumerate(features):
        col = i % 2
        row = i // 2
        left = Inches(0.8 + col * 6.0)
        top = Inches(1.65 + row * 1.15)
        panel = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, Inches(5.3), Inches(0.85)
        )
        panel.fill.solid()
        panel.fill.fore_color.rgb = PANEL_COLOR
        panel.line.color.rgb = BORDER
        tf = panel.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = feature
        p.font.name = "Arial"
        p.font.size = Pt(11.3)
        p.font.bold = True
        p.font.color.rgb = TITLE_COLOR
    add_slide_number(slide, 8)


def add_architecture_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, variant=2)
    add_logo(slide)
    add_title(slide, "Техникалық архитектура", "Платформа бірнеше өзара байланысты деңгейден тұрады")

    boxes = [
        ("Frontend", "HTML/CSS/JS\nЕкітілді интерфейс\nФормалар, карточкалар, кабинет"),
        ("Backend", "Django + DRF\nAPI, логика, аутентификация"),
        ("Database", "SQLite\nПайдаланушылар, өтініштер,\nжеңілдіктер, жаңалықтар"),
        ("Admin", "Jazzmin admin\nКонтент пен өтініштерді\nжедел басқару"),
    ]
    starts = [Inches(0.9), Inches(3.35), Inches(5.8), Inches(8.25)]
    shapes = []
    for (title, body), left in zip(boxes, starts):
        shp = add_panel(slide, left, Inches(2.45), Inches(2.1), Inches(2.2), title, body)
        shapes.append(shp)

    for i in range(len(shapes) - 1):
        x1 = shapes[i].left + shapes[i].width
        y1 = shapes[i].top + shapes[i].height // 2
        x2 = shapes[i + 1].left
        y2 = shapes[i + 1].top + shapes[i + 1].height // 2
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
        conn.line.color.rgb = ACCENT_3
        conn.line.width = Pt(2)

    add_panel(
        slide,
        Inches(1.3),
        Inches(5.2),
        Inches(10.3),
        Inches(0.95),
        "Технологиялық логика",
        "Пайдаланушы интерфейсі API арқылы сервермен байланысады, ал сервер деректер базасын және әкімшілік басқаруды қамтамасыз етеді.",
        fill_color=PANEL_DARK,
    )
    add_slide_number(slide, 9)


def add_metrics_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, variant=1)
    add_logo(slide)
    add_title(slide, "Прототиптің ағымдағы күйі", "Құжат дайындалған сәттегі жүйедегі нақты көрсеткіштер")

    chart_data = CategoryChartData()
    chart_data.categories = ["Санат", "Жеңілдік", "Жаңалық", "Өтініш", "Пайдаланушы"]
    chart_data.add_series("Саны", (5, 3, 7, 2, 7))

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.85),
        Inches(1.75),
        Inches(6.35),
        Inches(3.95),
        chart_data,
    ).chart
    chart.has_legend = False
    chart.value_axis.visible = False
    chart.category_axis.tick_labels.font.size = Pt(9)
    chart.category_axis.tick_labels.font.color.rgb = TITLE_COLOR
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = ACCENT
    chart.series[0].data_labels.show_value = True
    chart.series[0].data_labels.font.size = Pt(9)
    chart.series[0].data_labels.font.color.rgb = TITLE_COLOR
    chart.series[0].data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    add_stat_card(slide, Inches(7.55), Inches(1.85), Inches(2.0), Inches(1.15), "5", "Жеңілдік санаты")
    add_stat_card(slide, Inches(9.75), Inches(1.85), Inches(2.0), Inches(1.15), "3", "Жеңілдік карточкасы", color=ACCENT_3)
    add_stat_card(slide, Inches(7.55), Inches(3.2), Inches(2.0), Inches(1.15), "7", "Жаңалық", color=ACCENT_2)
    add_stat_card(slide, Inches(9.75), Inches(3.2), Inches(2.0), Inches(1.15), "2", "Тестік өтініш")
    add_stat_card(slide, Inches(8.65), Inches(4.55), Inches(2.0), Inches(1.15), "7", "Пайдаланушы", color=ACCENT_3)

    add_slide_number(slide, 10)


def add_swot_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, variant=2)
    add_logo(slide)
    add_title(slide, "SWOT-талдау", "Жобаның күшті және әлсіз жақтарын, мүмкіндіктері мен тәуекелдерін бағалау")

    boxes = [
        ("S — Strengths", ["Әлеуметтік маңыздылығы жоғары", "Жұмыс істейтін прототип", "Екітілді интерфейс", "Ақпарат + өтініш + сүйемелдеу бір жүйеде"]),
        ("W — Weaknesses", ["Жеңілдіктер базасы әлі толық емес", "AI толық енгізілмеген", "Тесттік база шағын", "Прототип деңгейіндегі инфрақұрылым"]),
        ("O — Opportunities", ["Аймақтар бойынша кеңейту", "ҮЕҰ және ЖОО-мен серіктестік", "AI-кеңесші қосу", "Мобильді нұсқа/қосымша"]),
        ("T — Threats", ["Масштабтауда қауіпсіздік талаптары", "Интеграция күрделілігі", "Ресми деректерді үнемі жаңарту қажеттігі", "Пайдаланушы сенімін сақтау маңызды"]),
    ]
    positions = [
        (Inches(0.75), Inches(1.6)),
        (Inches(6.65), Inches(1.6)),
        (Inches(0.75), Inches(4.0)),
        (Inches(6.65), Inches(4.0)),
    ]
    for (title, bullets), (left, top) in zip(boxes, positions):
        add_bullet_box(slide, left, top, Inches(5.3), Inches(1.95), title, bullets, fill_color=PANEL_DARK)
    add_slide_number(slide, 11)


def add_innovation_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, variant=1)
    add_logo(slide)
    add_title(slide, "Инновациясы және AI-перспективасы", "Rakhym технологияны әлеуметтік мәселені шешу құралы ретінде қолданады")

    add_bullet_box(
        slide,
        Inches(0.75),
        Inches(1.7),
        Inches(5.8),
        Inches(4.6),
        "Жобаның инновациясы",
        [
            "Әлеуметтік көмекті цифрландыру",
            "Пайдаланушы үшін бірыңғай кіру нүктесін ұсыну",
            "Ақпаратты, өтінішті және сүйемелдеуді біріктіру",
            "Еріктілер контурын платформа логикасына енгізу",
            "AI-модульдерге дайын архитектуралық негіз қалыптастыру",
        ],
    )
    add_bullet_box(
        slide,
        Inches(6.75),
        Inches(1.7),
        Inches(5.45),
        Inches(4.6),
        "Келесі AI-қадамдар",
        [
            "Пайдаланушы профиліне сай жеңілдіктерді ақылды іріктеу",
            "Типтік сұрақтар бойынша AI-кеңесші",
            "Өтініштерді автоматты маршруттау",
            "Өтініш шұғылдығын алдын ала бағалау",
            "Маманға кейс бойынша ұсыным беру",
        ],
        fill_color=PANEL_DARK,
    )
    add_slide_number(slide, 12)


def add_roadmap_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, variant=2)
    add_logo(slide)
    add_title(slide, "Шектеулері және даму жол картасы", "Прототиптің өсу нүктелері нақты көрініп тұр")

    add_bullet_box(
        slide,
        Inches(0.75),
        Inches(1.65),
        Inches(5.65),
        Inches(4.65),
        "Қазіргі шектеулер",
        [
            "Жеңілдіктер базасы толық көлемде толтырылмаған",
            "Жоба өнеркәсіптік мемлекеттік жүйе емес, прототип",
            "Кең көлемде енгізу үшін интеграция және қауіпсіздік күшейтілуі керек",
            "Тестілеу мен аналитиканы тереңдету қажет",
        ],
    )
    add_bullet_box(
        slide,
        Inches(6.65),
        Inches(1.65),
        Inches(5.6),
        Inches(4.65),
        "Келесі даму қадамдары",
        [
            "Жеңілдіктер мен жәрдемақылар базасын кеңейту",
            "Толыққанды AI-кеңесші енгізу",
            "Маман кабинеті мен аналитиканы күшейту",
            "Сыртқы сервистермен интеграция жасау",
            "Мобильді нұсқа немесе мобильді қосымша шығару",
        ],
        fill_color=PANEL_DARK,
    )
    add_slide_number(slide, 13)


def add_conclusion_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, variant=1)
    add_logo(slide, width=Inches(0.7))

    center = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(1.2), Inches(10.8), Inches(4.9)
    )
    center.fill.solid()
    center.fill.fore_color.rgb = PANEL_COLOR
    center.fill.transparency = 0.0
    center.line.color.rgb = BORDER

    tf = center.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = "ҚОРЫТЫНДЫ"
    r1.font.name = "Arial"
    r1.font.size = Pt(24)
    r1.font.bold = True
    r1.font.color.rgb = ACCENT

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = (
        "Rakhym — әлеуметтік көмекті\n"
        "түсініктірек, қолжетімдірек және адамға жақынырақ ететін цифрлық жоба."
    )
    r2.font.name = "Arial"
    r2.font.size = Pt(22)
    r2.font.bold = True
    r2.font.color.rgb = TITLE_COLOR

    p3 = tf.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run()
    r3.text = (
        "Бұл платформа ақпаратты, өтінішті, мәртебені бақылауды және\n"
        "сүйемелдеуді бір жүйеде біріктіріп, Social Tech бағытының нақты үлгісін көрсетеді."
    )
    r3.font.name = "Arial"
    r3.font.size = Pt(13)
    r3.font.color.rgb = TEXT_COLOR

    footer = slide.shapes.add_textbox(Inches(0.95), Inches(6.35), Inches(11.2), Inches(0.4))
    p = footer.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Назарларыңызға рақмет!"
    r.font.name = "Arial"
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = TITLE_COLOR
    add_slide_number(slide, 14)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def main() -> None:
    add_cover_slide()
    add_relevance_slide()
    add_problem_slide()
    add_goal_slide()
    add_audience_slide()
    add_solution_slide()
    add_journey_slide()
    add_features_slide()
    add_architecture_slide()
    add_metrics_slide()
    add_swot_slide()
    add_innovation_slide()
    add_roadmap_slide()
    add_conclusion_slide()

    prs.save(str(OUTPUT))
    print(OUTPUT)


if __name__ == "__main__":
    main()
